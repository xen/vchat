from __future__ import annotations

from io import BytesIO

import pytest
from docx import Document
from pptx import Presentation
from pypdf import PdfWriter
from scrapy.linkextractors import IGNORED_EXTENSIONS

from jobs.crawler.document_pipeline import (
    detect_binary_document_kind,
    extract_binary_url_document,
    extract_url_document,
)
from jobs.crawler.spiders.general import GeneralSpider


def test_extract_url_document_prefers_html_title_and_strips_auth_forms() -> None:
    html = """
    <html>
      <head>
        <title>Антибуллинговые карточки</title>
      </head>
      <body>
        <div class="auth-modal">
          <h2>Авторизация</h2>
          <form>
            <label>E-mail*</label>
            <input name="USER_LOGIN" />
            <label>Пароль*</label>
            <input type="password" name="USER_PASSWORD" />
          </form>
        </div>
        <main>
          <h1>Антибуллинговые карточки</h1>
          <p>Автор: АНО «БО "Журавлик"»</p>
          <p>Основное описание продукта.</p>
          <p>
            «Антибуллинговые карточки» помогают ребенку распознавать травлю,
            различать конфликт и буллинг, а также тренировать эмоциональную
            регуляцию, критическое мышление и безопасные способы реагирования
            на агрессию со стороны сверстников.
          </p>
          <p>
            Материал подходит для индивидуальной беседы и для групповой работы.
            Дополнительный текст нужен здесь только для того, чтобы тестовая
            страница не выглядела как пустая SPA-заглушка для эвристики
            `_static_text_length`.
          </p>
        </main>
      </body>
    </html>
    """

    content, title, meta = extract_url_document(
        "https://example.com/doc",
        html_body=html,
        content_type="text/html; charset=utf-8",
    )

    assert title == "Антибуллинговые карточки"
    assert "Авторизация" not in content
    assert "USER_LOGIN" not in content
    assert "Основное описание продукта." in content
    assert meta["extraction"]["extractor"] == "beautifulsoup"
    assert meta["extraction"]["boilerplate_removed_count"] >= 1


def test_extract_binary_url_document_extracts_docx_text_and_tables() -> None:
    document = Document()
    document.core_properties.title = "Metadata Support Rules"
    document.add_heading("Регламент поддержки", level=1)
    document.add_paragraph("PDF и Word документы должны попадать в базу знаний.")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Раздел"
    table.rows[0].cells[1].text = "Ответственный"
    table.rows[1].cells[0].text = "Индексация"
    table.rows[1].cells[1].text = "Команда поиска"
    buffer = BytesIO()
    document.save(buffer)

    content, title, meta = extract_binary_url_document(
        "https://example.com/files/Support%20Rules.docx",
        buffer.getvalue(),
        content_type=(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ),
    )

    assert title == "Metadata Support Rules"
    assert "PDF и Word документы должны попадать в базу знаний." in content
    assert "| Раздел | Ответственный |" in content
    assert meta["doc_type"] == "office"
    assert meta["extraction"]["extractor"] == "python-docx"
    assert meta["file"]["table_count"] == 1


def test_extract_binary_url_document_uses_query_filename_when_metadata_title_empty() -> (
    None
):
    document = Document()
    document.add_heading("Заголовок внутри документа", level=1)
    document.add_paragraph("Title должен быть взят из имени файла.")
    buffer = BytesIO()
    document.save(buffer)

    content, title, meta = extract_binary_url_document(
        "https://example.com/public/api/file?filename=consent.docx",
        buffer.getvalue(),
        content_type=(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ),
    )

    assert title == "consent"
    assert "Title должен быть взят из имени файла." in content
    assert meta["extraction"]["extractor"] == "python-docx"


def test_extract_binary_url_document_sniffs_docx_for_broad_binary_mime() -> None:
    document = Document()
    document.core_properties.title = "metadata-slug"
    document.add_paragraph("DOCX with a broad binary MIME still extracts.")
    buffer = BytesIO()
    document.save(buffer)

    content, title, meta = extract_binary_url_document(
        "https://example.com/public/api/file?filename=ignored.bin",
        buffer.getvalue(),
        content_type="application/octet-stream",
    )

    assert title == "metadata-slug"
    assert "DOCX with a broad binary MIME still extracts." in content
    assert meta["extraction"]["extractor"] == "python-docx"


def _sample_pptx_bytes() -> bytes:
    presentation = Presentation()
    presentation.core_properties.title = "Metadata Slide Deck"

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Презентация поддержки"
    title_slide.placeholders[1].text = "PowerPoint документы должны индексироваться."

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    table_slide.shapes.title.text = "Матрица ответственности"
    table = table_slide.shapes.add_table(2, 2, 0, 0, 4_000_000, 1_200_000).table
    table.cell(0, 0).text = "Раздел"
    table.cell(0, 1).text = "Ответственный"
    table.cell(1, 0).text = "Презентации"
    table.cell(1, 1).text = "Команда поиска"

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_extract_binary_url_document_extracts_pptx_slides_and_tables() -> None:
    content, title, meta = extract_binary_url_document(
        "https://example.com/files/support-deck.pptx",
        _sample_pptx_bytes(),
        content_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
    )

    assert title == "Metadata Slide Deck"
    assert "## Slide 1" in content
    assert "Презентация поддержки" in content
    assert "PowerPoint документы должны индексироваться." in content
    assert "| Раздел | Ответственный |" in content
    assert "| Презентации | Команда поиска |" in content
    assert meta["doc_type"] == "office"
    assert meta["extraction"]["extractor"] == "python-pptx"
    assert meta["file"]["slide_count"] == 2
    assert meta["file"]["table_count"] == 1


def test_extract_binary_url_document_sniffs_pptx_for_broad_binary_mime() -> None:
    content, title, meta = extract_binary_url_document(
        "https://example.com/public/api/file?filename=deck.bin",
        _sample_pptx_bytes(),
        content_type="application/octet-stream",
    )

    assert title == "Metadata Slide Deck"
    assert "PowerPoint документы должны индексироваться." in content
    assert meta["extraction"]["extractor"] == "python-pptx"


def test_extract_binary_url_document_accepts_pptx_with_legacy_powerpoint_mime() -> None:
    content, title, meta = extract_binary_url_document(
        "https://example.com/files/support-deck.pptx",
        _sample_pptx_bytes(),
        content_type="application/vnd.ms-powerpoint",
    )

    assert title == "Metadata Slide Deck"
    assert "PowerPoint документы должны индексироваться." in content
    assert meta["extraction"]["extractor"] == "python-pptx"


def test_extract_binary_url_document_rejects_legacy_ppt() -> None:
    with pytest.raises(ValueError, match=r"Legacy \.ppt files are not supported"):
        extract_binary_url_document(
            "https://example.com/files/legacy-deck.ppt",
            b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1legacy ppt bytes",
            content_type="application/vnd.ms-powerpoint",
        )


def test_extract_binary_url_document_returns_empty_content_for_image_only_pdf() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.add_metadata({"/Title": "Metadata Terms"})
    buffer = BytesIO()
    writer.write(buffer)

    content, title, meta = extract_binary_url_document(
        "https://example.com/files/Terms.pdf",
        buffer.getvalue(),
        content_type="application/pdf",
    )

    assert content == ""
    assert title == "Metadata Terms"
    assert meta["extraction"]["extractor"] == "pypdf"
    assert meta["extraction"]["fallback_used"] is False


def test_detect_binary_document_kind_uses_mime_before_url_path() -> None:
    assert (
        detect_binary_document_kind(
            b"<html><body>not a pdf</body></html>",
            "text/html",
        )
        is None
    )
    assert (
        detect_binary_document_kind(b"%PDF-1.7\n", "application/octet-stream") == "pdf"
    )


def test_general_spider_allows_word_pdf_and_pptx_links() -> None:
    spider = GeneralSpider(url="https://example.com", source_id=1, config={})
    try:
        assert {"pdf", "doc", "docx", "ppt", "pptx"}.issubset(IGNORED_EXTENSIONS)
        assert "pdf" not in spider._link_extractor.deny_extensions
        assert "doc" not in spider._link_extractor.deny_extensions
        assert "docx" not in spider._link_extractor.deny_extensions
        assert ".ppt" in spider._link_extractor.deny_extensions
        assert "pptx" not in spider._link_extractor.deny_extensions
    finally:
        spider.closed("test")
